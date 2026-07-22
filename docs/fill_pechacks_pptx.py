"""Fill the official PecHacks4.0 PPTX template with MediNexus AI content.

Replaces placeholder text run-by-run so all template fonts/styles/colors
are preserved exactly. Texts are sized to the template's boxes; a few
tight cells get a small font-size reduction or wrap-off (same fonts).
Output: Downloads/PecHacks4.0_MediNexus_AI.pptx
"""
from pptx import Presentation
from pptx.util import Pt, Inches

SRC = r"C:\Users\mshre\Downloads\PecHacks4.0 (1).pptx"
OUT = r"C:\Users\mshre\Downloads\PecHacks4.0_MediNexus_AI.pptx"

# (slide_no, shape_id) -> ordered texts for each meaningful run.
# Decoration runs ("|") are left untouched; leftover placeholder runs are blanked.
FILL = {
    # ---- Slide 1: Title ----
    (1, 12): ["TEAM MEDINEXUS"],
    (1, 13): ["MediNexus AI"],
    (1, 14): ["Healthcare"],
    (1, 15): [""],

    # ---- Slide 2: Problem Statement ----
    (2, 9): [
        "Rural clinics, small hospitals and disaster or low-connectivity zones work without "
        "AI-powered clinical decision support, because nearly every clinical AI tool today "
        "depends on the cloud — no internet means no intelligence. Patient records stay "
        "fragmented across paper files and disconnected spreadsheets, so deterioration goes "
        "unnoticed until it becomes an emergency, while overloaded doctors lose up to a third "
        "of their time to manual documentation. Sending sensitive patient data to third-party "
        "cloud servers adds privacy, compliance and cost barriers that many institutions cannot "
        "accept. There is no unified, offline-first platform combining records, live monitoring, "
        "medical knowledge and risk prediction on local, affordable hardware."
    ],

    # ---- Slide 3: Solution ----
    (3, 17): [
        "Solution",
        "MediNexus AI is an AI Clinical Copilot that runs entirely on a local machine — no "
        "internet, no cloud, no patient data ever leaving the premises. Five integrated "
        "modules, one intelligent AI Orchestrator:",
    ],
    (3, 18): ["MediNexus", " AI"],
    (3, 21): ["Patient records & live clinical dashboard"],
    (3, 22): ["Medical knowledge RAG via local LLM ",
              "Live vitals streaming & real-time alerts"],
    (3, 23): ["ML risk prediction & early warning"],
    (3, 24): ["Automated PDF/DOCX reports"],

    # ---- Slide 4: Dependencies & Showstoppers ----
    (4, 13): ["Local compute for the LLM — running offline AI needs decent hardware. "
              "Mitigated with quantized small models (Phi-3, Gemma 2B) on Ollama that run "
              "on a mid-range laptop, no GPU required."],
    (4, 14): ["Knowledge-base quality — RAG answers are only as good as the corpus. We use "
              "curated open medical guidelines (WHO, standard treatment protocols) and every "
              "answer cites its source."],
    (4, 15): ["LLM hallucination risk — a wrong answer is a patient-safety risk. Strict RAG "
              "grounding plus human-in-the-loop design: MediNexus assists clinicians, it "
              "never auto-diagnoses."],
    (4, 16): ["Real sensor integration — the demo uses a realistic vitals simulator over "
              "WebSockets. Production needs vendor APIs / HL7, and the streaming "
              "architecture is already built for that swap."],

    # ---- Slide 5: Additional Info ----
    (5, 8): [
        "MediNexus AI is offline-first by design: patient data never leaves the hospital "
        "premises, making it deployable where privacy rules, budgets or connectivity make "
        "cloud AI impossible. Around 65% of India's population is rural yet is served by "
        "under a third of its doctors — exactly where offline clinical intelligence "
        "matters most.",
        "A working Phase-1 build is already running: live clinical dashboard, patient CRUD "
        "with search, seeded data and risk analytics — remaining modules plug into the same "
        "FastAPI + React architecture. Below: rural population share (%), offline operation "
        "(%), documentation minutes saved per patient, laptops needed to deploy.",
    ],
    (5, 10): ["65"],
    (5, 11): ["100"],
    (5, 12): ["30"],
    (5, 13): ["1"],

    # ---- Slide 6: Tech Stack ----
    (6, 17): ["Frontend",
              "React + Tailwind CSS + Plotly — responsive clinical dashboard with "
              "interactive real-time charts."],
    (6, 19): ["Backend"],
    (6, 23): ["FastAPI + WebSockets — REST APIs for records and analytics, live streaming "
              "for vitals monitoring."],
    (6, 21): ["Database"],
    (6, 25): ["SQLite / PostgreSQL with SQLAlchemy — structured patients, vitals, reports "
              "and consultations."],
    (6, 16): ["AI / RAG",
              "FAISS + Sentence Transformers + LangChain — offline semantic search over "
              "curated medical knowledge."],
    (6, 18): ["Local LLM"],
    (6, 22): ["Ollama running Llama 3 / Phi-3 / Gemma — a fully offline clinical copilot "
              "with zero cloud calls."],
    (6, 20): ["ML & Reports"],
    (6, 24): ["PyTorch / TensorFlow risk models + ReportLab and python-docx for automated "
              "documentation."],

    # ---- Slide 7: Team ----
    (7, 10): ["Team Name: TEAM MEDINEXUS"],
    (7, 11): ["M. Shree"],
    (7, 12): ["Sanjay K"],
    (7, 13): ["Niranjan Saravanan"],
    (7, 14): ["Prajeet Joshua P"],
    (7, 15): ["Sanjeevan U S"],
    (7, 16): ["VIT Chennai"],
    (7, 17): ["VIT Chennai"],
    (7, 18): ["VIT Chennai"],
    (7, 19): ["Easwari Engineering College"],
    (7, 20): ["Sri Venkateswara College of Engg."],
    (7, 21): ["Chennai"],
    (7, 22): ["Chennai"],
    (7, 23): ["Chennai"],
    (7, 24): ["Chennai"],
    (7, 25): ["Sriperumbudur"],
    (7, 28): ["ML Engineer"],       # row 1 (M. Shree)
    (7, 29): ["Backend Dev"],       # row 2 (Sanjay K)
    (7, 30): ["Frontend Dev"],      # row 3 (Niranjan)
    (7, 26): ["Team Leader"],       # row 4 (Prajeet Joshua P)
    (7, 27): ["AI Engineer"],       # row 5 (Sanjeevan)
    (7, 31): ["—"], (7, 32): ["—"], (7, 33): ["—"], (7, 34): ["—"], (7, 35): ["—"],
}

# Minimal fit tweaks: same fonts, adjusted size / wrap / box width where the
# template box is too small for real content.
TWEAKS = {
    (1, 12): {"size": 28},                    # "TEAM MEDINEXUS" one line in 4.14"
    (1, 13): {"size": 30},                    # "MediNexus AI" one line in 3.05"
    (1, 14): {"size": 15},                    # domain + subtitle line in 7.5"
    (3, 18): {"size": 10},                    # "MediNexus / AI" inside 0.89" chip
    (3, 21): {"wrap": False},
    (3, 24): {"wrap": False},
    (7, 10): {"size": 20, "width": 5.5},      # team-name heading
    (7, 11): {"wrap": False}, (7, 12): {"wrap": False}, (7, 13): {"wrap": False},
    (7, 14): {"wrap": False}, (7, 15): {"wrap": False},
    (7, 16): {"wrap": False, "size": 10}, (7, 17): {"wrap": False, "size": 10},
    (7, 18): {"wrap": False, "size": 10}, (7, 19): {"wrap": False, "size": 10},
    (7, 20): {"wrap": False, "size": 9},
    (7, 21): {"wrap": False}, (7, 22): {"wrap": False}, (7, 23): {"wrap": False},
    (7, 24): {"wrap": False}, (7, 25): {"wrap": False},
    (7, 26): {"wrap": False, "size": 10}, (7, 27): {"wrap": False, "size": 10},
    (7, 29): {"wrap": False, "size": 10}, (7, 30): {"wrap": False, "size": 10},
}


def fill_shape(shape, texts):
    """Assign texts to meaningful runs in order; keep '|' decorations; blank leftovers."""
    queue = list(texts)
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if run.text.strip() in ("", "|"):
                continue  # decoration / spacing run — keep as-is
            run.text = queue.pop(0) if queue else ""


def break_before_second_pipe(shape):
    """Slide 3 shape 22 holds items 02 and 03 in one paragraph; add a hard
    line break before the second '|' so item 03 aligns with its number."""
    from pptx.oxml.ns import qn
    para = shape.text_frame.paragraphs[0]
    pipes = [r for r in para.runs if r.text.strip() == "|"]
    if len(pipes) >= 2:
        r2 = pipes[1]._r
        r2.addprevious(r2.makeelement(qn("a:br"), {}))


def tweak_shape(shape, opts):
    if "size" in opts:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                run.font.size = Pt(opts["size"])
    if opts.get("wrap") is False:
        shape.text_frame.word_wrap = False
    if "width" in opts:
        shape.width = Inches(opts["width"])


GH_URL = "https://github.com/MSHREE26092007/medinexus-ai"
SITE_URL = "https://mshree26092007.github.io/medinexus-ai/"


def add_links(prs):
    """Add clickable GitHub + live-demo links on slides 1 and 5.

    Hyperlink is set on the SHAPE (click action), not the text runs, so the
    text keeps the template-matching white Lato instead of the theme's dark
    hyperlink blue (unreadable on these backgrounds)."""
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    # slide idx: L, T, and per-link offset (dx, dy). Slide 1 stacks the links
    # left-aligned; slide 5 keeps them side by side.
    placements = {0: (0.55, 3.38, 0.0, 0.36), 4: (0.68, 4.78, 4.6, 0.0)}
    links = [("GitHub:  github.com/MSHREE26092007/medinexus-ai", GH_URL),
             ("Live demo:  mshree26092007.github.io/medinexus-ai", SITE_URL)]
    for idx, (left, top, dx, dy) in placements.items():
        slide = list(prs.slides)[idx]
        for i, (text, url) in enumerate(links):
            box = slide.shapes.add_textbox(Inches(left + i * dx), Inches(top + i * dy),
                                           Inches(4.4), Inches(0.35))
            tf = box.text_frame
            tf.word_wrap = False
            run = tf.paragraphs[0].add_run()
            run.text = text
            run.font.name = "Lato"
            run.font.size = Pt(13)
            run.font.bold = True
            run.font.underline = True
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            box.click_action.hyperlink.address = url


def restyle_title_slide(prs):
    """Slide 1: upper-left = domain, then project name, then small description."""
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    slide = list(prs.slides)[0]
    for sh in slide.shapes:
        if sh.shape_id == 13:          # "MediNexus AI" title: move up under domain
            sh.top = Inches(2.10)
    desc = slide.shapes.add_textbox(Inches(0.55), Inches(2.82), Inches(6.9), Inches(0.35))
    tf = desc.text_frame
    tf.word_wrap = False
    run = tf.paragraphs[0].add_run()
    run.text = ("An offline AI clinical copilot — unified records, RAG medical "
                "knowledge, live vitals & risk prediction.")
    run.font.name = "Lato"
    run.font.size = Pt(13)
    run.font.color.rgb = RGBColor(0xE8, 0xE8, 0xF2)


def add_flowchart(prs):
    """Slide 3: replace the idle gray placeholder graphics on the left with the
    solution flowchart image."""
    from pptx.util import Inches
    slide = list(prs.slides)[2]
    for sh in list(slide.shapes):
        if sh.shape_id in (5, 6, 7, 9, 11):   # gray placeholder blocks/mountains
            sh._element.getparent().remove(sh._element)
    slide.shapes.add_picture(r"C:\Users\mshre\medinexus-ai\docs\solution_flow.png",
                             Inches(0.13), Inches(2.90), Inches(3.85), Inches(2.55))


def main():
    prs = Presentation(SRC)
    done = set()
    for si, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            key = (si, shape.shape_id)
            if key in FILL and shape.has_text_frame:
                fill_shape(shape, FILL[key])
                if key in TWEAKS:
                    tweak_shape(shape, TWEAKS[key])
                if key == (3, 22):
                    break_before_second_pipe(shape)
                done.add(key)
    missing = set(FILL) - done
    if missing:
        raise SystemExit(f"ERROR — shapes not found: {sorted(missing)}")
    add_links(prs)
    restyle_title_slide(prs)
    add_flowchart(prs)
    prs.save(OUT)
    print(f"All {len(done)} shapes filled + links added. Saved: {OUT}")


if __name__ == "__main__":
    main()
